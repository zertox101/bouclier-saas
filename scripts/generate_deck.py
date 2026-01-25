"""Create the executive CyberDetect presentation into out/CyberDetect_Presentation.pptx."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

OUTPUT = Path(__file__).resolve().parent.parent / "out" / "CyberDetect_Presentation.pptx"
BACKGROUND_COLOR = RGBColor(6, 6, 9)
ACCENT_COLOR = RGBColor(33, 150, 243)
TEXT_COLOR = RGBColor(236, 239, 241)


def set_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BACKGROUND_COLOR


def add_textbox(slide, left, top, width, height, text, font_size=24, bold=False, color=TEXT_COLOR):
    box = slide.shapes.add_textbox(left, top, width, height)
    box.text_frame.clear()
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, height, title, bullets, notes):
    body_top = top
    if title:
        title_box = add_textbox(slide, left, top, width, Inches(0.8), title, font_size=28, bold=True)
        title_box.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR
        body_top = top + Inches(0.8)
    body = slide.shapes.add_textbox(left, body_top, width, height - (body_top - top))
    tf = body.text_frame
    tf.text = bullets[0]
    tf.paragraphs[0].level = 0
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.color.rgb = TEXT_COLOR
    for bullet in bullets[1:]:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.text = notes


def add_module_map(slide):
    add_textbox(slide, Inches(0.4), Inches(0.3), Inches(9), Inches(1), "Core Modules", font_size=30, bold=True)
    modules = [
        ("Surveillance", "Next.js dashboard & SSE charts"),
        ("Detection", "FastAPI rules + analytics + Ollama"),
        ("Telemetry", "Tools API, Kali traffic generator, Falco/auditd agents"),
        ("Support", "Postgres/Redis storage, Nginx redirector"),
    ]
    top = Inches(1.5)
    for name, desc in modules:
        add_textbox(slide, Inches(0.6), top, Inches(4), Inches(0.5), name, font_size=24, bold=True)
        add_textbox(slide, Inches(0.6), top + Inches(0.35), Inches(8), Inches(0.4), desc, font_size=18)
        top += Inches(1)
    add_textbox(slide, Inches(5.5), Inches(1.5), Inches(4.5), Inches(2.5), "Telemetry flows → FastAPI → Redis/Postgres → Frontend", font_size=20, bold=False)
    slide.notes_slide.notes_text_frame.text = "Explain the modular split: telemetry sources, detection, storage, and UI."


def add_story_slide(slide, title, map_items, notes):
    add_textbox(slide, Inches(0.4), Inches(0.3), Inches(9.2), Inches(0.8), title, font_size=28, bold=True)
    top = Inches(1.4)
    for idx, (heading, detail) in enumerate(map_items, start=1):
        add_textbox(slide, Inches(0.5), top, Inches(9), Inches(0.4), f"{idx}. {heading}", font_size=22, bold=True)
        add_textbox(slide, Inches(0.8), top + Inches(0.35), Inches(8.6), Inches(0.4), detail, font_size=18)
        top += Inches(0.9)
    slide.notes_slide.notes_text_frame.text = notes


def add_pipeline_slide(slide):
    add_textbox(slide, Inches(0.4), Inches(0.3), Inches(9.2), Inches(0.8), "Detection Pipeline", font_size=28, bold=True)
    pipeline = [
        "Telemetry ingestion (/api/events, /ws/traffic, Suricata/Falco/log forwarders).",
        "Storage in Postgres + Redis, plus enrichment via GeoIP + LLM context (Ollama).",
        "Analytics engine (Isolation Forest) and planned Sigma-style rule builder.",
        "Correlated alerts/incidents streamed back to UI and exported via incident packs.",
    ]
    add_bullets(slide, Inches(0.4), Inches(1), Inches(9), Inches(5), "", pipeline, "Highlight the mixed AI/rule strategy.")


def add_data_model_slide(slide):
    add_textbox(slide, Inches(0.4), Inches(0.3), Inches(9.2), Inches(0.8), "Data Model Highlights", font_size=28, bold=True)
    bullets = [
        "AlertEvent / EventLog: raw telemetry plus severity metadata.",
        "MlAlert / CorrelatedAlert: enriched with behavior scores and incidents.",
        "User / Sensor tables track onboarding state, heartbeats, and agent versions.",
        "Tools jobs and audit trail capture command, status, and output hashes.",
    ]
    add_bullets(slide, Inches(0.4), Inches(1.1), Inches(9.2), Inches(4.5), "", bullets, "Refer to backend/app/models/sql.py for schemas.")


def add_demo_slide(slide):
    add_textbox(slide, Inches(0.4), Inches(0.3), Inches(9.2), Inches(0.8), "Demo Ready Screenshots", font_size=28, bold=True)
    placeholders = [
        "KPI Dashboard / Traffic cards",
        "Alert Triage queue + timeline preview",
        "Incident pack export (PCAP snippet + hashes)",
    ]
    top = Inches(1.3)
    for text in placeholders:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.5), top, Inches(4.5), Inches(2.2)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(21, 32, 43)
        shape.line.color.rgb = ACCENT_COLOR
        textbox = shape.text_frame
        textbox.text = text
        textbox.paragraphs[0].font.size = Pt(18)
        textbox.paragraphs[0].font.color.rgb = TEXT_COLOR
        top += Inches(2.4)
    slide.notes_slide.notes_text_frame.text = "Point to live screenshots from the Next.js dashboard once deployed."


def add_roadmap_slide(slide):
    add_textbox(slide, Inches(0.4), Inches(0.3), Inches(9.2), Inches(0.8), "Roadmap (30/60/90)", font_size=28, bold=True)
    top = Inches(1.2)
    segments = [
        ("30 Days", "Wire new Sigma-style rule editor, improve API docs, ship overview refresh."),
        ("60 Days", "Ingest Suricata/Falco telemetry, add sensor health dashboards, enable JWT."),
        ("90 Days", "Launch incident pack export, ClickHouse/MinIO storage, and planner for new tools."),
    ]
    for label, desc in segments:
        add_textbox(slide, Inches(0.5), top, Inches(4), Inches(0.8), label, font_size=22, bold=True)
        add_textbox(slide, Inches(0.5), top + Inches(0.5), Inches(9), Inches(0.7), desc, font_size=18)
        top += Inches(1.2)
    slide.notes_slide.notes_text_frame.text = "Use this slide to align stakeholders on the next delivery horizon."


def add_security_slide(slide):
    add_textbox(slide, Inches(0.4), Inches(0.3), Inches(9.2), Inches(0.8), "Security & Privacy Considerations", font_size=28, bold=True)
    bullets = [
        "Docker isolation (Kali + tools API) plus cap_net_admin cap for packet inspection.",
        "JWT plans for backend plus Redis caching of telemetry to avoid data loss.",
        "Postgres encryption/backup policies and audit table for executed tools.",
        "Agent heartbeats detect tamper, while Falco/auditd watch for exec/write anomalies.",
    ]
    add_bullets(slide, Inches(0.4), Inches(1.1), Inches(9.2), Inches(4.5), "", bullets, "Reassure CISO that telemetry is contained in the SaaS boundary.")


def add_closing_slide(slide):
    add_textbox(slide, Inches(0.4), Inches(0.5), Inches(9.2), Inches(1), "Closing & Next Steps", font_size=32, bold=True)
    closing = [
        "CyberDetect is ready to showcase incident correlation and 1-click exports.",
        "Reach us at security@bouclier-saas.com for trial access, roadmap demos, or data reviews.",
    ]
    add_bullets(slide, Inches(0.4), Inches(1.5), Inches(9.2), Inches(4), "", closing, "End with thanks and contact info.")


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_textbox(slide, Inches(0.5), Inches(1.2), Inches(12), Inches(1), "CyberDetect", font_size=60, bold=True)
    add_textbox(
        slide,
        Inches(0.5),
        Inches(2.4),
        Inches(12),
        Inches(1),
        "Defense-grade telemetry, detection, and SOC workflows under one roof",
        font_size=28,
    )
    slide.notes_slide.notes_text_frame.text = "Introduce CyberDetect: telemetry, detection, and SOC in one SaaS."

    # Problem statement
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_bullets(
        slide,
        Inches(0.5),
        Inches(0.5),
        Inches(12),
        Inches(5),
        "Problem / Need",
        [
            "Security teams ingest too many siloed alerts without operator-friendly triage.",
            "Telemetry gaps, agent health blind spots, and missing correlation slow response.",
            "Existing SaaS lack integrated Purple Team validation, threat simulation, and exportable evidence.",
        ],
        "Frame the breakup: chaotic alerts, telemetry gaps, and lack of evidence export.",
    )

    # What CyberDetect does
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_bullets(
        slide,
        Inches(0.5),
        Inches(0.5),
        Inches(12),
        Inches(5),
        "What CyberDetect Does",
        [
            "Ingests IDS/agent telemetry (tools API, Kali generator, Falco/auditd, osquery) into FastAPI + Redis/Postgres.",
            "Applies ML/rule scoring with Ollama-powered context and canvases on the dashboard via SSE.",
            "Guides analysts through alert → incident → incident pack workflows with search + filters.",
            "Supports Purple Team validation, CSV exports, and on-prem tool execution orchestration.",
        ],
        "Quickly review product value: ingestion, detection, triage, and evidence.",
    )

    # Core modules slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_module_map(slide)

    # Main user journey
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    journey = [
        ("Alert discovery", "KPI cards + traffic metrics, filters, SSE updates drive triage queue."),
        ("Triage to incident", "Preview panel + timeline, assign/resolution actions, or federate to incident packs."),
        ("Evidence + reporting", "Export PCAP snippet, hashes, recommendations, and send to reports."),
    ]
    add_story_slide(slide, "Main User Journey", journey, "Walk through alert discovery, triage, and reporting.")

    # Detection pipeline
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_pipeline_slide(slide)

    # Data model highlight
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_data_model_slide(slide)

    # Demo placeholders
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_demo_slide(slide)

    # Roadmap
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_roadmap_slide(slide)

    # Security considerations
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_security_slide(slide)

    # Closing slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_closing_slide(slide)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(f"Wrote {OUTPUT}")


if __name__ == "__main__":
    build_presentation()
